from pptx import Presentation #pip install python-pptx
from pptx.dml.color import RGBColor
import sys
import os
import time
import xlwings as xw #pip install xlwings
from datetime import datetime
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

if getattr(sys, 'frozen', False):
    PATH_PASTA = os.path.dirname(sys.executable)
else:
    PATH_PASTA = os.path.dirname(__file__)

# PATH_PASTA = os.path.dirname(os.path.abspath(__file__))
print(PATH_PASTA)
PATH_TEMPLATE = PATH_PASTA +  "\\template_ncmr.pptx"
# print(PATH_TEMPLATE)
PATH_IMAGENS = PATH_PASTA + "\\arquivos\\"
# print(PATH_IMAGENS)
PATH_PLANILHA = PATH_PASTA + "\\banco_de_dados.xlsm"
# print(PATH_PLANILHA)

def left(string, n):
    """ Retorna a parte à esquerda de uma string, ou seja, o seu começo. O seu funcionamento é igual ao da função de mesmo nome do Excel.
     
    Parameters: 
    string (str): A string a ser cortada.
    n (int)     : O número de caracteres a serem incluídos na string resultado.

    Returns:
    str : Uma string de tamanho 'n' equivalente aos 'n' primeiros caracteres da variável 'string'. """
    return string[:n]

def right(string, n):
    """ Retorna a parte à direita de uma string, ou seja, os seus últimos caracteres. O seu funcionamento é igual ao da função de mesmo nome do Excel.
    
    Paramenters:
    string (str): A string a ser cortada.
    n (int)     : O número de caracteres a serem incluídos na string resultado.

    Returns:
    str : Uma string de tamanho 'n' equivalente aos 'n' últimos caracteres da variável 'string'. """
    posicao = len(string) - n
    return string[posicao:]

def adc_texto(paragrafo, conteudo, is_negrito):
    """ Adiciona um texto 'conteudo' a um certo 'paragrafo', em negrito ou não, a depender do valor de 'is_negrito'. 
    
    Parameters:
    paragrafo (obj. Paragraph): Um objeto do tipo Paragraph da biblioteca Python-pptx que irá conter o texto que se deseja inserir.
    conteudo (str): Uma string que contém o texto que se deseja inserir.
    is_negrito (bool): Variável que irá determinar a formatação do texto inserido. Se o valor for 'True', o texto será em negrito. Se o valor for 'False', não será negrito. """
    run = paragrafo.add_run() # Adiciona uma 'run'
    run.text = conteudo       # Insere texto nesta 'run'
    font = run.font           # Seleciona a fonte
    font.bold = is_negrito    # Altera a fonte para negrito ou não

def adc_campo(paragrafo, nome, valor):
    """ Adiciona um campo com o seu nome 'nome' e o seu valor 'valor' em um objeto do tipo 'Paragraph' chamado paragrafo.
    
    Parameters:
    paragrafo (obj. Paragraph): Um objeto do tipo Paragraph da biblioteca Python-pptx que irá conter o texto que se deseja inserir.
    nome (str): O nome do campo, que será inserido em negrito.
    valor (str): O valor de um determinado campo, que será inserido em formatação sem negrito. """
    adc_texto(paragrafo, nome, True)    # Adiciona um texto em negrito
    adc_texto(paragrafo, valor, False)  # Adiciona um texto sem negrito

def adc_shape(shape, campos, valores):
    """Preenche uma determinada 'shape' com os seus 'campos' e seus respectivos 'valores'.

    Parameters:
    shape (obj. Shape): Um objeto do tipo Shape da biblioteca Python-pptx que será preenchido e formatado.
    campos (list of str): Uma lista de strings que contém cada um dos campos que serão incluídos neste Shape.
    valores (list of str): Uma lista de strings que contém os valores de cada um dos campos que serão incluídos neste Shape. """

    txt_frame = shape.text_frame            # Seleciona o objeto que contém textos (text_frame) no objeto Shape. 
    txt_frame.clear()                       # Limpa esse objeto Text_frame
    p = txt_frame.paragraphs[0]             # Seleciona o primeiro parágrafo no Text_frame
    p.alignment = PP_ALIGN.JUSTIFY          # Alinhamento do texto será justificado
    for i in range(len(campos)):            # Adiciona os campos e os seus valores ao parágrafo p
        adc_campo(p, campos[i],valores[i])  

def add_slide_4block(presentation, master_number, layout_number):
    """ Cria um slide título no objeto 'presentation', de acordo com a escolha do slide master dada por 'master_number'.

    Parameters:
    presentation (obj. Presentation): Um objeto do tipo Presentation da biblioteca Python-pptx. Os slides mestres e os layouts selecionados dependem do template de que esta Presentation foi extraída.
    master_number (int): Um inteiro, que assume os valores de 0 ou 1. Com o valor 0 é selecionado o template em português, e com o valor 1 é selecionado o template em inglês. 
    layout_number (int): Um inteiro, equivalente ao layout que se deseja criar o slide. 0 para slide título, 1 para slide de imagens extra

    Returns:
    slide (obj. Slide): Retorna um objeto do tipo Slide da bilioteca Python-pptx. """
    master = presentation.slide_masters[master_number]              ## master contém o master de número master_number
    title_slide_layout = master.slide_layouts[layout_number]        ## Seleciona o slide de numero layout_number
    slide = presentation.slides.add_slide(title_slide_layout)       ## Cria um novo slide na configuracao do layout_number
    return slide

def four_block(presentation,dados,lingua="pt-br"):
    """ Cria, formata e popula um slide de 4Block com todas as informações necessárias.

    Parameters:
    presentation (obj. Presentation): Um objeto do tipo Presentation da biblioteca Python-pptx. Os slides mestres e os layouts selecionados dependem do template de que esta Presentation foi extraída.
    dados (dict): Um dicionário que contém todos os dados de texto que serão incluídos no Four-Block.
    lingua (str): Uma string com a língua do slide do diagrama que se deseja. Os seus valores padrões são 'pt-br' ou 'English'. """
    ## Seleciona o template a ser usado (ingles ou portugues)
    if (lingua == "pt-br"):
        master_number = 0
        campos_descricao = ["Qual é o problema? ", "\n\nQual é o estado correto? "]
        campos_dados_relevantes = ["Part Number: ", 
                    "\nSerial Number: ",
                    "\nQuantidade: ",
                    "\nModelo da Loco: ",
                    "\nNúmero da Loco: ",
                    "\nFornecedor / Fabricante: "]
        campos_disposicao = ["O que fazer para resolver o problema? ","\nLista de materiais e horas necessárias: ", "\nCusto estimado do retrabalho (materiais + horas): ", "\nResponsável local por este report e data: "]
        valores_disposicao = ["N/A", "Não aplicável durante a emissão de documentos.", "R$ 00", "Input de "+ dados["email_responsavel"] + " em " + dados["data_requisicao"]+"."]
        titulos_figuras = ["Correto (estado esperado)", "Incorreto (estado encontrado)"]
        campo_prioridade = ["Prioridade: "]
    else:
        master_number = 1 
        campos_descricao = ["What is the problem? ", "\n\nWhat is the correct form/behavior? "]
        campos_dados_relevantes = ["Part Number: ", 
                    "\nSerial Number: ",
                    "\nQuantity: ",
                    "\nLoco model: ",
                    "\nLoco number: ",
                    "\nSupplier / Manufacturer: "]
        campos_disposicao = ["What to do to solve the problem? ","\nList of necessary materials and hours: ", "\nEstimated cost of reworking (materials + hours): ", "\nPerson responsible for this report and date: "]
        valores_disposicao = ["N/A", "Not applicable during the emission of documents.", "R$ 00", "Input of "+ dados["email_responsavel"] + " in " + dados["data_requisicao"]+"."]
        titulos_figuras = ["Correct (specified condition)", "Incorrect (found condition)"]
        campo_prioridade = ["Priority: "]
    slide = add_slide_4block(presentation,master_number,0)

    ##########################################################################
    ## Inserir título e código da NCMR
    titulo = slide.shapes[0]
    titulo.text = dados["cod"] + " - PN " + dados['part_number'] + " - " + dados["descricao_part"]
    
    ##########################################################################
    ## Inserir código da NCMR
    cod_ncmr = slide.shapes[1]
    cod_ncmr.text = dados["cod"]    

    ##########################################################################
    ## Inserir descrição
    descricao = slide.shapes[2]
    try:
        if (dados["Informativa"]=="Sim"):
            texto_descricao = dados["problema"] + "\n\n" + "NCMR INFORMATIVA"
    except:
        texto_descricao = dados["problema"]
    
    try:
        valores_descricao = [texto_descricao, dados["correto"]]
    except:
        valores_descricao = [texto_descricao, ""]
    adc_shape(descricao,campos_descricao, valores_descricao)

    ##########################################################################
    ## Inserir prioridade
    prioridade = slide.shapes[9]
    try:
        valor_prioridade = [dados["prioridade"]]
        adc_shape(prioridade,campo_prioridade, valor_prioridade)
        cor = cor_prioridade(dados["prioridade"])
        prioridade.fill.solid()
        prioridade.fill.fore_color.rgb  = cor
    except:
        pass

    ##########################################################################
    ## Inserir dados relevantes
    dados_relevantes = slide.shapes[3]
    try:
        valores_dados_relevantes = [dados["part_number"],
                      dados["serial_number"],
                      dados["qtd_itens"],
                      dados["modelo_locomotiva"],
                      dados["numero_locomotiva"],
                      dados["fornecedor"] + " / " + dados["fabricante"]] 
    except:
        valores_dados_relevantes = [dados["part_number"],
                      dados["serial_number"],
                      "1",
                      dados["modelo_locomotiva"],
                      dados["numero_locomotiva"],
                      "" + " / " + ""] 
    adc_shape(dados_relevantes, campos_dados_relevantes, valores_dados_relevantes)

    ##########################################################################
    ## Inserir disposição
    disposicao = slide.shapes[4]
    adc_shape(disposicao, campos_disposicao,valores_disposicao)
    
    ##########################################################################
    ## Inserir figuras
    diretorio_imagens = PATH_IMAGENS+dados["PowerAppsId"]
    imagens(presentation, slide, diretorio_imagens, master_number, titulos_figuras)

def imagens(presentation, slide, diretorio_imagens, master_number, titulos_figuras):
    """Insere as imagens no Four-Block.
    
    Parameters:
    presentation (obj. Presentation): Um objeto do tipo Presentation da biblioteca Python-pptx. Os slides mestres e os layouts selecionados dependem do template de que esta Presentation foi extraída.
    slide (obj. Slide): Um objeto do tipo Slide da bilioteca Python-pptx, é o slide título onde as primeiras figuras serão inseridas.
    diretorio_imagens (string): Uma string que contém o endereço do diretório com as imagens a ser inseridas.
    master_number (int): Um inteiro, que assume os valores de 0 ou 1. Com o valor 0 é selecionado o template em português, e com o valor 1 é selecionado o template em inglês. 
    titulo_figuras (array): Um array de strings com os possíveis títulos para identificar as figuras como Correta ou Incorreta. """
    tem_correta = False
    for name in os.listdir(diretorio_imagens): # Itera sobre os arquivos das figuras e testa se há figuras da forma correta.
        if name.startswith("Correta"):
            tem_correta = True
            break
    #print(tem_correta)
    if tem_correta:                            # Se tem figura da forma correta, executa este código 
        legenda_figura(slide.shapes[7],titulos_figuras[0],RGBColor(71,212,90)) # Adiciona legenda com a cor verde, indicando figura Correta
        legenda_figura(slide.shapes[8],titulos_figuras[1],RGBColor(255,0,0))   # Adiciona legenda com a cor vermelha, indicando figura Incorreta
        for name in os.listdir(diretorio_imagens):              # Itera sobre os arquivos
            path_imagem = diretorio_imagens+"/"+name            
            if name.endswith(".pptx"):                          # Se for um arquivo de formato .pptx, não fazer nada
                pass
            elif name.startswith("Correta"):                    
                if name.startswith("Correta(1)"):               # Se for o arquivo com nome Correta(1), insere no primeiro placeholder do slide título
                    figura = slide.shapes[5]                    
                    inserir_imagem(figura, path_imagem)
                else:                                           # Se for o arquivo com nome Correta(n), insere num slide extra 
                    slide_img_extra(presentation, master_number, path_imagem, titulos_figuras[0], RGBColor(71,212,90))
            elif name.startswith("Incorreta"):                  # Se for o arquivo com nome Incorreta(1), insere no segundo placeholder do slide título
                if name.startswith("Incorreta(1)"):
                    figura = slide.shapes[6]        
                    inserir_imagem(figura, path_imagem)
                else:                                          # Se for o arquivo com nome Incorreta(n), insere num slide extra 
                    slide_img_extra(presentation, master_number, path_imagem, titulos_figuras[1], RGBColor(255,0,0))
    
    else:                                       # Se não tem figura da forma correta, executa este código
        legenda_figura(slide.shapes[7],titulos_figuras[1],RGBColor(255,0,0)) # Adiciona legenda com a cor vermelha, indicando figura Incorreta
        legenda_figura(slide.shapes[8],titulos_figuras[1],RGBColor(255,0,0)) # Adiciona legenda com a cor vermelha, indicando figura Incorreta
        for name in os.listdir(diretorio_imagens):
            path_imagem=diretorio_imagens+"/"+name
            if name.endswith(".pptx"):                          # Se for um arquivo de formato .pptx, não fazer nada
                pass
            elif name.startswith("Incorreta(1)"):               # Se for o arquivo com nome Incorreta(1), insere no segundo placeholder do slide título
                #print(name)
                figura = slide.shapes[5]
                inserir_imagem(figura, path_imagem)
            elif name.startswith("Incorreta(2)"):               # Se for o arquivo com nome Incorreta(2), insere no segundo placeholder do slide título
                #print(name)
                figura = slide.shapes[6]
                inserir_imagem(figura, path_imagem)
            else:                                               # Se for o arquivo com nome Incorreta(1), insere num slide extra 
                slide_img_extra(presentation, master_number, path_imagem ,titulos_figuras[1], RGBColor(255,0,0))

def legenda_figura(placeholder, texto, cor):
    """ Insere as legendas das figuras com o texto selecionado e a cor selecionada.

    Parameters:
    placeholder (obj. Shape): Um objeto do tipo Shape da biblioteca Python-pptx que será preenchido e formatado.
    texto (string): Um string com texto de legenda da figura.
    cor (RGBColor): Um objeto RGBColor da biblioteca Python-pptx importado por pptx.dml.color """
    placeholder.text = texto
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb  = cor

def slide_img_extra(presentation, master_number, path_imagem, legenda, cor):
    """ Insere um slide de imagem extra e insere a imagem.
    
    Parameters:
    presentation (obj. Presentation): Um objeto do tipo Presentation da biblioteca Python-pptx. Os slides mestres e os layouts selecionados dependem do template de que esta Presentation foi extraída.
    master_number (int): Um inteiro, que assume os valores de 0 ou 1. Com o valor 0 é selecionado o template em português, e com o valor 1 é selecionado o template em inglês. 
    path_imagem (string): Uma string com o endereço completo da imagem a ser adicionada.
    legenda (string): Um string com texto de legenda da figura.
    cor (RGBColor): Um objeto RGBColor da biblioteca Python-pptx importado por pptx.dml.color
    """
    # criar novo slide
    master = presentation.slide_masters[master_number]
    slide_layout = master.slide_layouts[1]       ## Seleciona o slide de numero layout_number
    novo_slide = presentation.slides.add_slide(slide_layout)
    # inserir legenda
    legenda_figura(novo_slide.shapes[0],legenda,cor)
    # inserir figura
    figura = novo_slide.shapes[1]
    inserir_imagem(figura, path_imagem)

def get_header_row(tabela):
    """ Extrai os nomes dos campos do cabeçalho de uma variável 'tabela' do tipo Table (da bibliotec xlwings) e retorna um dicionário 'relacao' em que as chaves são os nomes dos campos e os seus valores são os endereços das respectivas colunas na tabela do Excel.

    Parameters:
    tabela (Table): Um objeto do tipo Table da biblioteca xlwings, que equivale a uma tabela extraída do Excel.

    Returns:
    relacao (dict) : Um dicionário em que as chaves são os nomes dos campos do cabeçalho da Table 'tabela' e os seus valores são os endereços das respectivas colunas  na tabela do Excel. 

    """
    endereco_cabecalho = tabela.header_row_range
    relacao={}
    for i in endereco_cabecalho:
        coluna = left((i.address),2)
        relacao[i.value]=coluna

    return relacao

def get_data(sheet, tabela, posicao):
    """
    Extrai os dados da Sheet 'sheet', da Table 'tabela' e da linha da posição de uma célula no Excel representada pela string 'posicao'. Retorna um dicionario 'dados' com as chaves sendo os nomes dos campos na Table 'tabela' e os valores sendo os valores de cada campo na linha de 'posicao'.

    Parameters:
    sheet (obj. Sheet): Um objeto do tipo Sheet da biblioteca xlwings, que equivale a uma planilha extraída do Excel.
    tabela (obj. Table): Um objeto do tipo Table da biblioteca xlwings, que equivale a uma tabela extraída do Excel.
    posiscao (str): Uma string que contém o endereço da célula selecionada no Excel.
    
    Returns:
    dados (dict): Um dicionário em que as chaves são os nomes dos campos na Table 'tabela' e os valores são os valores de cada campo na linha de 'posicao'.
    """
    # Imprime na tela a posição da célula selecionada no Excel.
    print("\tPosição:", posicao)
    # Imprime na tela o valor célula selecionada no Excel.
    print("\tCódigo: ", sheet[posicao].value, "\n")

    relacao = get_header_row(tabela)
    
    # Cria um dicionário vazio
    dados={}
    # A linha de que os dados serão extraídos.
    linha = right(posicao,2)
    
    # Itera nas chaves (nomes dos campos) do dicionário 'relacao'
    for i in relacao:
        coluna = relacao[i]
        endereco = coluna + linha

        # Transformando todos os valores da planilha em strings:
        if type(sheet[endereco].value) == datetime:
            dados[i]=(sheet[endereco].value).strftime('%d/%m/%Y')
        elif type(sheet[endereco].value) == float:
            dados[i]=str(int(sheet[endereco].value))
        elif sheet[endereco].value == None:
            dados[i]=" "
        else:
            dados[i]=sheet[endereco].value
        #print(i,": ",dados[i])
    dados["descricao_part"] = dados["descricao_part"].replace("/","-")
    return dados

def inserir_imagem(figura, path_imagem):
    """ Tenta inserir a imagem e retorna uma mensagem de erro se não for possível.
    
    Parameters:
    figura (obj. Shape): Um objeto do tipo Shape da biblioteca Python-pptx, onde a figura será inserida.
    path_imagem (string): Uma string com o endereço completo da imagem a ser adicionada.
    """
    try:
        figura.insert_picture(path_imagem)
    except:
        print("Não foi possível inserir a imagem com o código, insira manualmente.") 

def cor_prioridade(texto):
    switch = {
        "1. Emergente":RGBColor(240,0,0),
        "2. Muito urgente": RGBColor(255,166,0),
        "3. Urgente": RGBColor(250,240,0),
        "4. Pouco urgente": RGBColor(85,185,0),
        "5. Não urgente": RGBColor(0,200,255)
    }
    return switch.get(texto, RGBColor(255,255,255))

########## Main ##########
n = len(sys.argv)
#### Arguments passed ####
if n > 1:
    print("Programa executado a partir do Excel.")
    aba=sys.argv[1]
    posicao=sys.argv[2]
else:
    print("Programa executado sem interação do Excel. Será acessada a primeira linha da tabela.")
    aba='ncmr'
    posicao="$A$8"

print("Acessando a planilha...")
wb = xw.Book(PATH_PLANILHA)
sheet = wb.sheets[aba]
tabela = sheet.tables(aba)
print("Extraindo as informações...")
print("\tTabela:", aba)
dados = get_data(sheet, tabela, posicao)
path_4block = PATH_IMAGENS+dados["PowerAppsId"].strip()+"\\"
salvar = True
for file in os.listdir(path_4block):
    if file.startswith(dados["cod"]) and file.endswith(".pptx"):
        salvar = False
        print("Four-Block já existe.")
        break
      
if salvar:
    print("Criando o Four-Block...")
    prs = Presentation(PATH_TEMPLATE)
    four_block(prs, dados)
    four_block(prs, dados, "English")
    #time.sleep(5)
  
    if dados["cod"] == " ":
        prs.save(path_4block + "NCMR-RC-0000-000000 - PN " + dados["part_number"] + ' - ' + dados["descricao_part"] + ".pptx")
        print("\tFour-Block criado.")
        print("\tAbrindo o Four-Block...")
        os.startfile(path_4block + "NCMR-RC-0000-000000 - PN " + dados["part_number"] + ' - ' + dados["descricao_part"] + ".pptx")
    else:
        prs.save(path_4block + dados["cod"] + " - PN " + dados["part_number"] + ' - ' + dados["descricao_part"] + ".pptx")
        print("\tFour-Block criado.")
        print("\tAbrindo o Four-Block...")
        os.startfile(path_4block + dados["cod"] + " - PN " + dados["part_number"] + ' - ' + dados["descricao_part"] + ".pptx")
else:
    print("\tAbrindo o Four-Block...")
    os.startfile(path_4block + dados["cod"] + " - PN " + dados["part_number"] + ' - ' + dados["descricao_part"] + ".pptx")
time.sleep(5)


num_int = 3